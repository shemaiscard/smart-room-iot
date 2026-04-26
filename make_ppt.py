from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Professional Color Palette ─────────────────────────────────
BG_COLOR       = RGBColor(0xFF, 0xFF, 0xFF)   # Pure White
TEXT_PRIMARY   = RGBColor(0x1A, 0x1A, 0x1A)   # Near Black
TEXT_SECONDARY = RGBColor(0x4A, 0x4A, 0x4A)   # Dark Gray
ACCENT_BLUE    = RGBColor(0x00, 0x56, 0xD2)   # Corporate Blue
ACCENT_GRAY    = RGBColor(0xE0, 0xE0, 0xE0)   # Light Gray divider

def set_bg(slide, color=BG_COLOR):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, font_size=18, bold=False,
                 color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)

    # Title
    add_text_box(slide, title_text, 0.8, 0.5, 11.7, 0.8,
                 font_size=36, bold=True, color=ACCENT_BLUE)

    # Divider line
    add_rect(slide, 0.8, 1.4, 11.7, 0.02, ACCENT_GRAY)

    # Bullets
    y = 1.8
    for bullet in bullets:
        if bullet.startswith("__sub__"):
            txt = "      " + bullet[7:]
            fs, clr, bi = 18, TEXT_SECONDARY, False
        elif bullet.startswith("__h__"):
            txt = bullet[5:]
            fs, clr, bi = 22, TEXT_PRIMARY, True
        else:
            txt = "-  " + bullet
            fs, clr, bi = 20, TEXT_PRIMARY, False
        
        add_text_box(slide, txt, 0.8, y, 11.7, 0.5, font_size=fs, bold=bi, color=clr)
        y += 0.6
    return slide

# ═══════════════════════════════════════════════════════════════
# Slide 1 – Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# Blue accent block at the bottom
add_rect(slide, 0, 6.5, 13.33, 1.0, ACCENT_BLUE)

# Main Title
add_text_box(slide, "Smart Room", 1.0, 1.8, 11.3, 1.2,
             font_size=54, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)
add_text_box(slide, "IoT-Based Environmental Monitoring and Automated Control System",
             1.0, 3.0, 11.3, 0.8, font_size=24, bold=False,
             color=TEXT_SECONDARY, align=PP_ALIGN.LEFT)

# Footer info
add_text_box(slide, "Internet of Things (IoT) | Midterm Project",
             1.0, 6.75, 8.0, 0.5, font_size=16, color=RGBColor(255, 255, 255),
             align=PP_ALIGN.LEFT)
add_text_box(slide, "Shema Nkindi Giscard",
             9.0, 6.75, 3.3, 0.5, font_size=16, color=RGBColor(255, 255, 255),
             align=PP_ALIGN.RIGHT, bold=True)

# ═══════════════════════════════════════════════════════════════
# Slide 2 – Problem Definition
# ═══════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Problem Definition", [
    "Lack of real-time environmental monitoring in traditional living spaces",
    "High energy wastage due to unmanaged lighting and cooling systems",
    "Absence of remote monitoring and control capabilities for users",
    "Undetected environmental discomfort impacting productivity and health",
    "Potential safety hazards such as gas leaks or unauthorized access",
    "__h__Proposed Solution",
    "__sub__A specialized Smart Room prototype focusing on integrated automation",
    "__sub__Real-time monitoring and control using high-efficiency IoT protocols",
])

# ═══════════════════════════════════════════════════════════════
# Slide 3 – Hardware Components
# ═══════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Technologies Used - Hardware", [
    "ESP32 Microcontroller: Central processing with integrated WiFi connectivity",
    "DHT11 / DHT22 Sensors: Precision temperature and humidity measurement",
    "PIR Sensor: Advanced infrared technology for motion detection",
    "LDR: Sensing ambient light levels for automated illumination",
    "5V Relay Module: Safe switching of electrical appliances",
    "SG90 Servo Motor: Electromechanical control for smart entry systems",
    "__h__Network Protocol",
    "__sub__MQTT Protocol over 2.4GHz Wireless Network",
])

# ═══════════════════════════════════════════════════════════════
# Slide 4 – Software Architecture
# ═══════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Technologies Used - Software", [
    "Arduino Framework: Hardware-level firmware for data acquisition",
    "MQTT Broker: Cloud-based communication hub for data routing",
    "Python 3: Server-side logic and simulation processing",
    "Streamlit Framework: Dynamic web application for data visualization",
    "Plotly Library: Interactive high-resolution analytical charting",
    "__h__Deployment and Repository",
    "__sub__Hosted on Streamlit Cloud with direct GitHub integration",
    "__sub__Source Code: github.com/shemaiscard/smart-room-iot",
])

# ═══════════════════════════════════════════════════════════════
# Slide 5 – Evidence of Progress
# ═══════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Evidence of Progress", [
    "Successful assembly and testing of the ESP32 hardware prototype",
    "Validated MQTT data transmission cycles between device and cloud broker",
    "Fully operational web dashboard displaying real-time sensor metrics",
    "Implementation of remote control logic for lights, fans, and door systems",
    "Verified seamless data flow from sensors to end-user interface",
    "Development of simulation tools for comprehensive system validation",
])

# ═══════════════════════════════════════════════════════════════
# Slide 6 – Challenges and Strategic Steps
# ═══════════════════════════════════════════════════════════════
add_bullet_slide(prs, "Challenges and Next Steps", [
    "__h__Technical Challenges",
    "Managing connectivity stability in varied network environments",
    "Optimization of power consumption for extended device lifecycle",
    "Implementation of robust security layers for remote data access",
    "__h__Future Enhancements",
    "Integration of voice-based command interfaces",
    "Migration to industrial-grade cloud database solutions",
    "Refinement of hardware design for a more compact form factor",
])

# ═══════════════════════════════════════════════════════════════
# Slide 7 – Conclusion
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.1, ACCENT_BLUE)

add_text_box(slide, "Conclusion", 1.0, 2.5, 11.3, 1.2,
             font_size=48, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Thank you for your attention", 1.0, 3.8, 11.3, 0.8,
             font_size=24, bold=False, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
add_text_box(slide, "Source Code Available at github.com/shemaiscard/smart-room-iot",
             1.0, 5.0, 11.3, 0.5, font_size=18, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)

output = "Smart_Room_IoT_Project_v3.pptx"
prs.save(output)
print(f"Successfully saved: {output}")
